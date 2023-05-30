import re
from pptx import Presentation

class PowerPointReader:
    def __init__(self, file_path):
        self.presentation = Presentation(file_path)

    def clean_text(self, text):
        # Remove leading/trailing whitespace and unnecessary characters
        return re.sub(r'\s+', ' ', text).strip()

    def prepare_prompts(self):
        """
        Parses the PowerPoint presentation and returns a list of slide contents
        after cleaning the text.
        """
        for slide in self.presentation.slides:
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            slide_text += run.text + "\n"
            yield self.clean_text(slide_text)


